"""Extract raw text from uploaded files, and chunk it for retrieval.

Text-based formats only (no images): txt/md/csv/tsv/json/code read directly; pdf/docx/pptx
go through lazily-imported parsers (installed on the PC). Big files get chunked and fed to
the retriever; small ones can be injected whole. This is the "attach a file" feature.
"""
from __future__ import annotations

import os

PLAIN = {"txt", "md", "markdown", "csv", "tsv", "json", "py", "js", "ts", "html", "htm",
         "css", "c", "cpp", "java", "go", "rs", "rb", "sh", "yaml", "yml", "toml", "tex", "log"}


def extract_text(path: str) -> str:
    ext = path.lower().rsplit(".", 1)[-1] if "." in path else ""
    if ext in PLAIN:
        with open(path, encoding="utf-8", errors="replace") as f:
            return f.read()
    if ext == "pdf":
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n\n".join((p.extract_text() or "") for p in pdf.pages)
    if ext == "docx":
        import docx
        return "\n".join(p.text for p in docx.Document(path).paragraphs)
    if ext == "pptx":
        from pptx import Presentation
        out = []
        for slide in Presentation(path).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return "\n".join(out)
    raise ValueError(f"unsupported attachment type: .{ext} (text-based files only)")


def chunk(text: str, size: int = 1000, overlap: int = 120) -> list[str]:
    """Split into overlapping character windows, preferring paragraph boundaries."""
    text = text.strip()
    if len(text) <= size:
        return [text] if text else []
    chunks, i = [], 0
    while i < len(text):
        end = min(i + size, len(text))
        if end < len(text):
            nl = text.rfind("\n", i + size // 2, end)   # try to break on a newline
            if nl != -1:
                end = nl
        chunks.append(text[i:end].strip())
        if end >= len(text):
            break
        i = max(end - overlap, i + 1)
    return [c for c in chunks if c]


def extract_and_chunk(path: str, size: int = 1000, overlap: int = 120) -> list[str]:
    return chunk(extract_text(path), size, overlap)
